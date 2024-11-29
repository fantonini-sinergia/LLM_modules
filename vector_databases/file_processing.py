from pptx import Presentation
from docx import Document
import pandas as pd
import fitz
import PyPDF2
# # easyocr
# import easyocr

# tesseract
import pytesseract
from PIL import Image
pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'

def read_file(file):
    """
    input: 
        - a dict with args:
            -'name': name of the file with extension
            -'path': the path including the file name and extension
    return:
        - a dict with args (if file correctly read)
            -'text'
            -'file_extension'
            -'pages_start_char'
        - 0 (otherwise)
    """

    print(f"\n    Reading {file['name']} ...")

    readable = True
    text = ""
    pages_start_char = [0]
    file_extension = file['name'].split('.')[-1]

    # read txt
    if file_extension.upper() == 'TXT':
        """
        only get text with .txt
        """
        try:
            with open(file['path'], 'r') as f:
                text = f.read()
        except:
            readable = False
    
    # read words
    elif file_extension.upper() in ['DOC', 'DOCX']:
        """
        only get text with .doc
        """
        try:
            doc = Document(file['path'])
            text = "\n".join([para.text for para in doc.paragraphs])
        except:
            readable = False

    # read excels
    elif file_extension.upper() in ['XLSX', 'XLS']:
        """
        get text and page start char with excels
        """
        try:
            dfs = pd.read_excel(file['path'], sheet_name=None, header=None)
            for i in dfs.keys():
                dfs[i] = dfs[i].dropna(axis=0, how='all')
                dfs[i] = dfs[i].dropna(axis=1, how='all')
                for index, row in dfs[i].iterrows():
                    if not row.isnull().any():
                        dfs[i] = dfs[i].loc[index:].reset_index(drop=True)
                        break
                new_column_names = dfs[i].iloc[0]
                dfs[i] = dfs[i][1:]
                dfs[i].columns = new_column_names
                dfs[i].reset_index(drop=True, inplace=True)
                text_per_page = dfs[i].to_markdown()
                text += text_per_page
                pages_start_char.append(len(text_per_page) + pages_start_char[-1])
            pages_start_char.pop()
        except:
            readable = False

    # read powerpoints
    elif file_extension.upper() in ['PPT', 'PPTX']:
        """
        get text and page start char with .ppt
        """
        try:
            prs = Presentation(file['path'])
            for slide in prs.slides:
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text_per_page = run.text
                            text += text_per_page
                pages_start_char.append(len(text_per_page) + pages_start_char[-1])
            pages_start_char.pop()
        except:
            readable = False
    
    # read PDFs
    elif file_extension.upper() == 'PDF':
        """
        get text and page start char with .pdf
        """
        try:
            doc = fitz.open(file['path'])

            # # easyocr
            # reader = easyocr.Reader(['it','en']) # this needs to run only once to load the model into memory
            
            for page in doc:
                if len(page.get_text("dict")["blocks"]) < 1:

                    # # easyocr
                    # pix = page.get_pixmap().tobytes()
                    # text_per_page = "".join(reader.readtext(pix, detail = 0))

                    # tesseract
                    pix = page.get_pixmap()
                    image = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                    text_per_page = pytesseract.image_to_string(image,lang='ita')

                else:
                    text_per_page = page.get_text()
                text += text_per_page
                pages_start_char.append(len(text_per_page) + pages_start_char[-1])
            pages_start_char.pop()
            doc.close()
        except:
            readable = False

    # unsupported extension
    else:
        readable = False
    
    if readable:
        print(f"    {len(text)} characters read\n")
        return {
            "text": text, 
            "pages_start_char": pages_start_char, 
            "file_extension": file_extension
            }
    else:
        print(f"    0 characters read\n")
        return {
            "text": '', 
            "pages_start_char": 0, 
            "file_extension": None
            }


def extract_page(file_path, page_number):
    print(f"\n    Extracting page {page_number} from {file_path}...")

    with open(file_path, 'rb') as input_pdf_file:
        reader = PyPDF2.PdfReader(input_pdf_file)
        writer = PyPDF2.PdfWriter()

        page = reader.pages[page_number]
        writer.add_page(page)

        file_and_page_name = f"temporary_pdf_pag{page_number}.pdf"

        with open(file_and_page_name, 'wb') as output_pdf_file:
            writer.write(output_pdf_file)
    
    print(f"    page extracted\n")
    return file_and_page_name